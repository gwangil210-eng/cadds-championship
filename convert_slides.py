"""
MP4 발표 영상 → 제출용 PPTX 변환 스크립트
5팀 CADDS 발표 슬라이드 최종안
"""
import subprocess
import sys
import shutil
from pathlib import Path

from pptx import Presentation
from pptx.util import Cm

VIDEO_PATH = Path("/root/.claude/uploads/7b5791b7-a0ec-4953-8d3d-5837bc0698de/1b5910e5-_____12.mp4")
FRAMES_DIR = Path("/tmp/cadds_frames")
OUTPUT_PPTX = Path("/home/user/cadds-championship/05팀_발표슬라이드_최종안.pptx")

SLIDE_WIDTH_CM = 33.867
SLIDE_HEIGHT_CM = 19.05


def reset_frames_dir():
    if FRAMES_DIR.exists():
        shutil.rmtree(FRAMES_DIR)
    FRAMES_DIR.mkdir(parents=True)


def extract_frames_scene_change(threshold: float) -> list[Path]:
    cmd = [
        "ffmpeg", "-i", str(VIDEO_PATH),
        "-vf", f"select=gt(scene\\,{threshold})",
        "-vsync", "vfr",
        "-q:v", "2",
        str(FRAMES_DIR / "frame_%04d.jpg"),
        "-y",
    ]
    result = subprocess.run(cmd, capture_output=True, text=True)
    if result.returncode != 0:
        print(f"ffmpeg 경고: {result.stderr[-300:]}", file=sys.stderr)
    return sorted(FRAMES_DIR.glob("frame_*.jpg"))


def extract_frames_fixed_fps(fps: float) -> list[Path]:
    cmd = [
        "ffmpeg", "-i", str(VIDEO_PATH),
        "-vf", f"fps={fps}",
        "-q:v", "2",
        str(FRAMES_DIR / "frame_%04d.jpg"),
        "-y",
    ]
    subprocess.run(cmd, capture_output=True, text=True)
    return sorted(FRAMES_DIR.glob("frame_*.jpg"))


def extract_frames() -> list[Path]:
    for threshold in (0.25, 0.35, 0.45):
        reset_frames_dir()
        print(f"장면 변화 감지 추출 중 (임계값={threshold})...")
        frames = extract_frames_scene_change(threshold)
        count = len(frames)
        print(f"  추출된 프레임: {count}개")
        if 3 <= count <= 60:
            return frames
        if count > 60:
            continue

    print("장면 감지 실패 → 1fps 고정 추출로 폴백...")
    reset_frames_dir()
    frames = extract_frames_fixed_fps(1.0)
    print(f"  추출된 프레임: {len(frames)}개")
    return frames


def build_pptx(frames: list[Path]) -> None:
    prs = Presentation()
    prs.slide_width = Cm(SLIDE_WIDTH_CM)
    prs.slide_height = Cm(SLIDE_HEIGHT_CM)

    blank_layout = prs.slide_layouts[6]

    for i, img_path in enumerate(frames, 1):
        slide = prs.slides.add_slide(blank_layout)
        slide.shapes.add_picture(
            str(img_path),
            left=0,
            top=0,
            width=prs.slide_width,
            height=prs.slide_height,
        )
        if i % 5 == 0 or i == len(frames):
            print(f"  슬라이드 {i}/{len(frames)} 완료")

    prs.save(str(OUTPUT_PPTX))


def main():
    if not VIDEO_PATH.exists():
        sys.exit(f"오류: 영상 파일을 찾을 수 없습니다 → {VIDEO_PATH}")

    print(f"소스: {VIDEO_PATH}")
    print(f"출력: {OUTPUT_PPTX}")
    print()

    frames = extract_frames()
    if not frames:
        sys.exit("오류: 프레임을 추출하지 못했습니다.")

    print(f"\nPPTX 생성 중 ({len(frames)}개 슬라이드)...")
    build_pptx(frames)

    size_mb = OUTPUT_PPTX.stat().st_size / 1024 / 1024
    print(f"\n완료: {OUTPUT_PPTX.name} ({size_mb:.1f} MB, {len(frames)}슬라이드)")


if __name__ == "__main__":
    main()
