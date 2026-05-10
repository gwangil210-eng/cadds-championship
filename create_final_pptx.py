#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
create_final_pptx.py
걷기 실천율과 만성질환 상관관계 분석 모델 — Project Team 5
15장 제출형 발표 PPTX 생성
"""

from pathlib import Path
from pptx import Presentation
from pptx.util import Cm, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from lxml import etree

OUTPUT = Path("/home/user/cadds-championship/05팀_발표슬라이드_최종안.pptx")

# ── Palette ────────────────────────────────────────────────────────────────────
TEAL    = RGBColor(0x1D, 0x7B, 0x7B)
NAVY    = RGBColor(0x1A, 0x2B, 0x4A)
CRIMSON = RGBColor(0x85, 0x1A, 0x20)
BLUE_C  = RGBColor(0x1A, 0x40, 0x80)
AMBER   = RGBColor(0xA0, 0x60, 0x20)
GREEN   = RGBColor(0x2D, 0x8B, 0x6B)
LGRAY   = RGBColor(0xF2, 0xF2, 0xF7)
MGRAY   = RGBColor(0x88, 0x88, 0x99)
TEXT    = RGBColor(0x1E, 0x1E, 0x32)
BORDER  = RGBColor(0xCC, 0xCC, 0xDD)
WHITE   = RGBColor(0xFF, 0xFF, 0xFF)
RED     = RGBColor(0xCC, 0x22, 0x33)
DKBG    = RGBColor(0x1A, 0x2A, 0x40)

KR   = "맑은 고딕"
CODE = "Consolas"

W = Cm(33.867)
H = Cm(19.05)


# ── Core helpers ───────────────────────────────────────────────────────────────

def new_prs():
    prs = Presentation()
    prs.slide_width  = W
    prs.slide_height = H
    return prs


def blank(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


def rect(slide, x, y, w, h, fill=None, line_color=None, line_pt=0.75):
    from pptx.enum.shapes import MSO_SHAPE_TYPE
    shp = slide.shapes.add_shape(1, x, y, w, h)
    shp.line.fill.background()
    if fill:
        shp.fill.solid()
        shp.fill.fore_color.rgb = fill
    else:
        shp.fill.background()
    if line_color:
        shp.line.color.rgb = line_color
        shp.line.width = Pt(line_pt)
    else:
        shp.line.fill.background()
    return shp


def txt(slide, text, x, y, w, h,
        size=12, bold=False, color=None, align=PP_ALIGN.LEFT,
        font=None, wrap=True, italic=False, space_before=0):
    txb = slide.shapes.add_textbox(x, y, w, h)
    tf  = txb.text_frame
    tf.word_wrap = wrap
    p   = tf.paragraphs[0]
    p.alignment = align
    if space_before:
        p.space_before = Pt(space_before)
    run = p.add_run()
    run.text = text
    run.font.name  = font or KR
    run.font.size  = Pt(size)
    run.font.bold  = bold
    run.font.italic = italic
    run.font.color.rgb = color or TEXT
    return txb


def multiline_txt(slide, lines, x, y, w, h, size=11, color=None, leading=1.15):
    """lines: list of (text, bold, color_override)"""
    txb = slide.shapes.add_textbox(x, y, w, h)
    tf  = txb.text_frame
    tf.word_wrap = True
    first = True
    for item in lines:
        if isinstance(item, str):
            t, b, c = item, False, None
        else:
            t, b, c = item[0], item[1] if len(item) > 1 else False, item[2] if len(item) > 2 else None
        if first:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        run = p.add_run()
        run.text = t
        run.font.name  = KR
        run.font.size  = Pt(size)
        run.font.bold  = b
        run.font.color.rgb = c or color or TEXT
    return txb


def top_bar(slide):
    rect(slide, Cm(0), Cm(0), Cm(26), Cm(0.17), fill=TEAL)
    rect(slide, Cm(26), Cm(0), W - Cm(26), Cm(0.17), fill=NAVY)


def page_num(slide, n):
    txt(slide, f"{n:02d}", Cm(31.2), Cm(18.2), Cm(2.4), Cm(0.7),
        size=11, color=MGRAY, align=PP_ALIGN.RIGHT)


def title_bar(slide, title, y=Cm(0.8)):
    rect(slide, Cm(1.3), y, Cm(0.22), Cm(0.72), fill=TEAL)
    txt(slide, title, Cm(1.7), y - Cm(0.04), Cm(30), Cm(0.8),
        size=21, bold=True, color=NAVY)


def subtitle_line(slide, text, y=Cm(1.85)):
    txt(slide, text, Cm(1.3), y, Cm(31.2), Cm(0.65), size=10.5, color=TEXT)


def card(slide, x, y, w, h, border_col, label, title, body,
         lsz=7.5, tsz=13, bsz=9.5):
    rect(slide, x, y, w, Cm(0.15), fill=border_col)
    bg = rect(slide, x, y + Cm(0.15), w, h - Cm(0.15), fill=None, line_color=BORDER, line_pt=0.5)
    txt(slide, label, x + Cm(0.28), y + Cm(0.3),  w - Cm(0.4), Cm(0.45), size=lsz, color=MGRAY)
    txt(slide, title, x + Cm(0.28), y + Cm(0.72), w - Cm(0.4), Cm(0.65), size=tsz, bold=True, color=TEXT)
    txt(slide, body,  x + Cm(0.28), y + Cm(1.42), w - Cm(0.4), h - Cm(1.6),  size=bsz, color=TEXT, wrap=True)


def icon_circle(slide, label, x, y, r=Cm(0.45)):
    rect(slide, x - r, y - r, r*2, r*2, fill=LGRAY)
    txt(slide, label, x - r, y - r*0.55, r*2, r, size=9, color=TEAL,
        align=PP_ALIGN.CENTER, bold=True)


# ── Slide 01 — Title ───────────────────────────────────────────────────────────

def s01_title(prs):
    sl = blank(prs)
    # Top gradient bar
    rect(sl, Cm(0), Cm(0), W, Cm(0.22), fill=TEAL)
    rect(sl, Cm(24), Cm(0), W - Cm(24), Cm(0.22), fill=NAVY)

    # Logo circle
    rect(sl, Cm(14.43), Cm(0.8), Cm(4.5), Cm(4.5), fill=None, line_color=TEAL, line_pt=2)
    txt(sl, "◑", Cm(14.43), Cm(1.2), Cm(4.5), Cm(3.5),
        size=44, color=TEAL, align=PP_ALIGN.CENTER, font="Segoe UI Symbol")

    # KDT label
    txt(sl, "KDT AI 전문가 과정", Cm(8), Cm(5.6), Cm(17), Cm(0.7),
        size=12, color=TEAL, align=PP_ALIGN.CENTER, bold=False)

    # Main title
    txt(sl, "걷기 실천율과 만성질환", Cm(4), Cm(6.5), Cm(26), Cm(1.8),
        size=42, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
    txt(sl, "상관관계 분석 모델", Cm(4), Cm(8.1), Cm(26), Cm(1.8),
        size=42, bold=True, color=TEAL, align=PP_ALIGN.CENTER)

    # Subtitle
    txt(sl, "전국 229개 시군구 건강 데이터 기반 다층 통계 분석 보고서",
        Cm(4), Cm(10.2), Cm(26), Cm(0.7),
        size=13, color=MGRAY, align=PP_ALIGN.CENTER)

    # Team badges
    badge_data = [("Project Team 5", NAVY, WHITE, Cm(10.0)),
                  ("송광일", LGRAY, TEXT, Cm(14.3)),
                  ("김정범", LGRAY, TEXT, Cm(17.1)),
                  ("박수빈", LGRAY, TEXT, Cm(19.9)),
                  ("조정연", LGRAY, TEXT, Cm(22.7))]
    for label, bg_c, fg_c, bx in badge_data:
        rect(sl, bx, Cm(11.4), Cm(2.5), Cm(0.75), fill=bg_c, line_color=BORDER, line_pt=0.3)
        txt(sl, label, bx + Cm(0.1), Cm(11.42), Cm(2.3), Cm(0.7),
            size=10, bold=(bg_c == NAVY), color=fg_c, align=PP_ALIGN.CENTER)

    # Bottom stats (3 columns)
    stats = [
        ("분석 대상", "전국 229개 시군구 패널", Cm(4.5)),
        ("핵심 방법론", "다층 통계 & 머신러닝", Cm(13.4)),
        ("연구 목표", "데이터 기반 보건 정책 제안", Cm(22.3)),
    ]
    rect(sl, Cm(1.3), Cm(16.9), Cm(31.2), Cm(0.03), fill=LGRAY)
    for label, val, sx in stats:
        txt(sl, label, sx, Cm(17.0), Cm(8.5), Cm(0.5),
            size=8, color=MGRAY)
        txt(sl, val,   sx, Cm(17.55), Cm(8.5), Cm(0.7),
            size=12, bold=True, color=TEXT)

    page_num(sl, 1)
    return sl


# ── Slide 02 — 연구 배경 및 추진 전략 ────────────────────────────────────────────

def s02_background(prs):
    sl = blank(prs)
    top_bar(sl)
    title_bar(sl, "연구 배경 및 추진 전략")
    subtitle_line(sl, "From Local to National:  국소적인 표본의 한계를 극복하고 전국 단위로 분석 범위를 확장합니다.")

    cards = [
        (Cm(1.3),  Cm(3.0), CRIMSON, "AS-IS",   "기존 연구의 한계",
         "기존 대전광역시 단위(n = 30)의 인프라 중심 분석은 표본 부족으로 인해 통계적 유의성 확보가 불가능한 한계가 존재했습니다."),
        (Cm(17.5), Cm(3.0), BLUE_C,  "TO-BE",   "분석 범위 전국 확장",
         "통계적 신뢰도 확보를 위해 전국 229개 시군구 패널 데이터로 조사 대상을 전수 확장하였습니다."),
        (Cm(1.3),  Cm(9.5), TEAL,    "FOCUS",   "주제 고도화 (행태 중심)",
         "물리적 '인프라' 관점에서 벗어나, 시민들의 실제 '건강 행태(걷기 실천율)'와 질환 지표 간의 직접적 상관관계를 규명합니다."),
        (Cm(17.5), Cm(9.5), AMBER,   "METHOD",  "정밀 상관 모델 구축",
         "지역별 고유 환경 특성이 통제된 다층적 데이터 구조를 반영하여, 데이터 기반의 명확한 의사결정 모델을 제안합니다."),
    ]
    for (cx, cy, bc, lbl, ttl, body) in cards:
        card(sl, cx, cy, Cm(15.0), Cm(5.8), bc, lbl, ttl, body)

    page_num(sl, 2)
    return sl


# ── Slide 03 — 목차 ────────────────────────────────────────────────────────────

def s03_toc(prs):
    sl = blank(prs)
    top_bar(sl)
    title_bar(sl, "목차  /  Table of Contents")

    sections = [
        ("01", "표지 및 연구 개요",    "걷기 실천율과 만성질환 상관관계 분석 모델 소개"),
        ("02", "연구 배경 및 추진 전략", "From Local to National — 분석 범위 전국 확장"),
        ("04", "분석 데이터 개요",      "질병관리청 지역사회건강조사(2008-2025) 기반 변수 정의"),
        ("05", "분석 아키텍처",         "R 생태계 기반 4단계 파이프라인 구조"),
        ("06", "데이터 전처리 파이프라인", "숫자형 변환 → 결측치 제거 → 공간 결합 → 정합성 검증"),
        ("07", "10대 건강 지표 상관관계 히트맵", "걷기율↔비만율 r=−0.52, 비만율↔고혈압 r=0.64"),
        ("08", "다중 회귀 분석: 비만율 영향 요인", "생활습관 변수 통제 후 독립적 영향력 산출"),
        ("09", "만성질환 영향 분석 (간접 경로 규명)", "걷기 → 비만 감소 → 만성질환 감소 메커니즘"),
        ("10", "K-Means 군집 분석 (K=3)", "건강 위험군 / 중간군 / 양호군 지역 분류"),
        ("11", "머신러닝 비만 예측 변수 중요도", "Random Forest 변수 기여도 분석"),
        ("12", "패널 회귀 분석 (Fixed Effects)", "지역 고유 특성 통제 후 걷기율 인과 효과 입증"),
        ("13", "정책 시뮬레이션 (2026-2030 예측)", "걷기율 +1~5%p 시 비만율 3.5%p 억제 효과"),
        ("14", "4대 건강지표 통합 미래 트렌드",   "2030년 비만·당뇨·고혈압 거시적 변화 전망"),
        ("15", "최종 분석 결론 및 정책 제언",     "유의성 검증, 경로 규명, 지역 격차, 복합 정책 제안"),
    ]

    col1 = sections[:7]
    col2 = sections[7:]

    for col_idx, col_items in enumerate([col1, col2]):
        cx = Cm(1.3) if col_idx == 0 else Cm(17.5)
        for row_idx, (num, ttl, desc) in enumerate(col_items):
            cy = Cm(2.5) + row_idx * Cm(2.2)
            rect(sl, cx, cy, Cm(1.0), Cm(0.6), fill=TEAL)
            txt(sl, num, cx, cy + Cm(0.05), Cm(1.0), Cm(0.55),
                size=9, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
            txt(sl, ttl,  cx + Cm(1.15), cy,          Cm(14.0), Cm(0.6),
                size=10.5, bold=True, color=NAVY)
            txt(sl, desc, cx + Cm(1.15), cy + Cm(0.6), Cm(14.0), Cm(0.55),
                size=8.5, color=MGRAY)

    page_num(sl, 3)
    return sl


# ── Slide 04 — 분석 데이터 개요 ──────────────────────────────────────────────────

def s04_data(prs):
    sl = blank(prs)
    top_bar(sl)
    title_bar(sl, "분석 데이터 개요")
    subtitle_line(sl,
        "질병관리청 지역사회건강조사(2008-2025) 기반  전국 229개 시군구, 17개년 10대 건강행태 변수 통합")

    # Table
    rows_data = [
        ("구분",           "변수명",              "산출 설명 (지표 기준)", True),
        ("독립 변수 (행태)", "걷기_실천율",
         "최근 1주일간 1일 30분 이상 걷기를 주 5일 이상 실천한 비율", False),
        ("",              "비만율",
         "체질량지수(BMI) 25 이상인 사람의 비율", False),
        ("종속 변수 (결과)", "고혈압_진단율",
         "의사로부터 고혈압 진단을 받은 적이 있는 사람의 비율", False),
        ("",              "당뇨_진단율",
         "의사로부터 당뇨 진단을 받은 적이 있는 사람의 비율", False),
        ("통제 변수 (위험)", "흡연/음주/스트레스",
         "현재 흡연자 비율, 고위험 음주 행태 비율, 스트레스 인지 수준", False),
        ("공간 변수",       "SHP (Shapefile)",
         "전국 시군구 단위 행정구역 공간 데이터 결합", False),
    ]

    ty = Cm(2.85)
    col_ws = [Cm(7.5), Cm(7.0), Cm(17.5)]
    col_xs = [Cm(1.3), Cm(8.8), Cm(15.8)]
    row_h  = Cm(1.6)

    for ri, (c0, c1, c2, is_header) in enumerate(rows_data):
        ry = ty + ri * row_h
        bg = LGRAY if is_header else (WHITE if ri % 2 == 0 else RGBColor(0xFA, 0xFA, 0xFF))
        for ci, (cx, cw, ct) in enumerate(zip(col_xs, col_ws, [c0, c1, c2])):
            fill = bg if not is_header else RGBColor(0xE4, 0xEF, 0xEF)
            rect(sl, cx, ry, cw, row_h, fill=fill, line_color=BORDER, line_pt=0.4)
            bold_cell = is_header or (ci == 0 and c0 != "")
            txt(sl, ct, cx + Cm(0.25), ry + Cm(0.35), cw - Cm(0.3), row_h - Cm(0.4),
                size=10.5 if not is_header else 11, bold=bold_cell,
                color=NAVY if is_header else (TEXT if ci > 0 else NAVY))

    page_num(sl, 4)
    return sl


# ── Slide 05 — 분석 아키텍처 ──────────────────────────────────────────────────────

def s05_architecture(prs):
    sl = blank(prs)
    top_bar(sl)
    title_bar(sl, "분석 아키텍처 (4단계 파이프라인)")
    subtitle_line(sl,
        "데이터 수집부터 정책 해석까지 R 생태계 기반 단일 파이프라인으로 연결하여 정합성을 확보했습니다.")

    steps = [
        ("①", "1. 데이터 수집",    "readr, sf 패키지",
         "17개년 건강조사 CSV 및 행정구역 SHP 공간 데이터 통합 로드"),
        ("②", "2. 데이터 전처리",  "dplyr, tidyr 패키지",
         "결측치 및 변수 정제, 공간/속성 데이터 JoinKey 기반 병합"),
        ("③", "3. 분석 엔진",      "상관/회귀/군집",
         "Random Forest, 패널 회귀(Fixed Effects) 등 다층적 통계 분석"),
        ("④", "4. 시각화",         "ggplot2, patchwork",
         "통합 대시보드 및 지리정보 기반 인터랙티브 시각화 구현"),
    ]
    colors = [TEAL, CRIMSON, BLUE_C, GREEN]

    sw = Cm(7.5)
    for i, (icon, title, pkg, desc) in enumerate(steps):
        sx = Cm(1.3) + i * (sw + Cm(0.5))
        sy = Cm(4.5)
        rect(sl, sx, sy, sw, Cm(0.2), fill=colors[i])
        # Icon circle
        rect(sl, sx + Cm(3.0), sy + Cm(0.5), Cm(1.5), Cm(1.5),
             fill=RGBColor(0xEE, 0xF6, 0xF6), line_color=colors[i], line_pt=1)
        txt(sl, icon, sx + Cm(3.0), sy + Cm(0.6), Cm(1.5), Cm(1.3),
            size=18, bold=True, color=colors[i], align=PP_ALIGN.CENTER)
        # Step title
        txt(sl, title, sx, sy + Cm(2.3), sw, Cm(0.75),
            size=13, bold=True, color=TEXT, align=PP_ALIGN.CENTER)
        # Package
        txt(sl, pkg, sx, sy + Cm(3.1), sw, Cm(0.55),
            size=10, color=colors[i], align=PP_ALIGN.CENTER, font=CODE)
        # Description
        txt(sl, desc, sx + Cm(0.2), sy + Cm(3.75), sw - Cm(0.4), Cm(3.5),
            size=9.5, color=TEXT, align=PP_ALIGN.CENTER, wrap=True)

    page_num(sl, 5)
    return sl


# ── Slide 06 — 데이터 전처리 파이프라인 ──────────────────────────────────────────

def s06_preprocessing(prs):
    sl = blank(prs)
    top_bar(sl)
    title_bar(sl, "데이터 전처리 파이프라인")
    subtitle_line(sl,
        "분석 결과의 신뢰도를 결정하는 핵심 단계로, 철저한 검증을 통해 분석용 마트를 구축했습니다.")

    steps = [
        (NAVY,    "1. 숫자형 변환",  "as.numeric()",
         "문자열로 수집된 로우 데이터를 연산 가능한 수치 데이터로 일괄 변환"),
        (NAVY,    "2. 결측치 제거",  "drop_na()",
         "핵심 분석 변수 기준 결측치가 존재하는 불완전한 레코드를 데이터셋에서 제외"),
        (NAVY,    "3. 공간 결합",    "left_join()",
         "공통 JoinKey 생성 후 SHP 공간 속성과 10대 건강 지표 데이터 매칭"),
        (BLUE_C,  "4. 정합성 검증",  "anti_join()",
         "행정구역 통폐합에 따른 매칭 오류를 사전 탐지하여 데이터 신뢰성 100% 확보"),
    ]

    for i, (bg, title, fn, desc) in enumerate(steps):
        sy = Cm(3.2) + i * Cm(3.45)
        rect(sl, Cm(1.3), sy, Cm(5.8), Cm(2.7), fill=bg)
        txt(sl, title, Cm(1.3), sy + Cm(0.9), Cm(5.8), Cm(0.85),
            size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

        # Function tag
        rect(sl, Cm(28.5), sy + Cm(0.95), Cm(4.0), Cm(0.72),
             fill=LGRAY, line_color=BORDER, line_pt=0.5)
        txt(sl, fn, Cm(28.5), sy + Cm(0.98), Cm(4.0), Cm(0.65),
            size=10.5, color=TEAL, font=CODE, align=PP_ALIGN.CENTER)

        # Description
        txt(sl, desc, Cm(7.8), sy + Cm(0.9), Cm(20.5), Cm(0.9),
            size=11, color=TEXT)

    page_num(sl, 6)
    return sl


# ── Slide 07 — 10대 건강 지표 상관관계 히트맵 ───────────────────────────────────

def s07_heatmap(prs):
    sl = blank(prs)
    top_bar(sl)
    title_bar(sl, "10대 건강 지표 상관관계 히트맵")

    # Correlation matrix
    headers = ["변수", "걷기율", "비만율", "고혈압", "당뇨"]
    rows = [
        ["걷기율",  "1.00", "−0.52", "−0.28", "−0.31"],
        ["비만율",  "−0.52", "1.00",  " 0.64",  " 0.58"],
        ["고혈압",  "−0.28", " 0.64", "1.00",  " 0.42"],
        ["당뇨",    "−0.31", " 0.58", " 0.42", "1.00"],
    ]
    special = {(1,2): RED, (1,3): RED, (2,1): RED, (2,2): GREEN, (2,3): GREEN}
    # special cells: row 1 (비만율 row) col 2,3 = red. row 2 (고혈압) col 1 = green, etc.

    tx0 = Cm(1.3)
    ty0 = Cm(2.3)
    cws = [Cm(3.2), Cm(3.2), Cm(3.2), Cm(3.2), Cm(3.2)]
    rh  = Cm(1.35)

    for ci, h in enumerate(headers):
        cx = tx0 + ci * Cm(3.2)
        rect(sl, cx, ty0, Cm(3.2), rh, fill=RGBColor(0xE0, 0xED, 0xED))
        txt(sl, h, cx + Cm(0.1), ty0 + Cm(0.35), Cm(3.0), rh - Cm(0.4),
            size=11, bold=True, color=NAVY, align=PP_ALIGN.CENTER)

    for ri, row in enumerate(rows):
        for ci, val in enumerate(row):
            cx = tx0 + ci * Cm(3.2)
            cy = ty0 + (ri+1) * rh
            is_diag = (ri == ci - 1)  # diagonal cells
            bg = RGBColor(0xEE, 0xF5, 0xF5) if is_diag else (LGRAY if ri % 2 == 0 else WHITE)
            rect(sl, cx, cy, Cm(3.2), rh, fill=bg, line_color=BORDER, line_pt=0.3)
            # Color highlight for key correlations
            c = TEXT
            if ri == 0 and ci == 2:  c = GREEN   # 걷기↔비만 -0.52
            elif ri == 1 and ci in (3,4): c = RED  # 비만↔고혈압/당뇨
            elif ri == 2 and ci == 2: c = RED
            elif ri == 3 and ci == 2: c = RED
            bold_val = (c != TEXT)
            txt(sl, val, cx + Cm(0.1), cy + Cm(0.35), Cm(3.0), rh - Cm(0.4),
                size=12, bold=bold_val, color=c, align=PP_ALIGN.CENTER)

    # Key findings
    fx = Cm(17.5)
    # Finding 1
    rect(sl, fx, Cm(2.3), Cm(15.5), Cm(5.5), fill=None,
         line_color=GREEN, line_pt=1.5)
    txt(sl, "r = −0.52  (음의 상관)", fx + Cm(0.3), Cm(2.55), Cm(15.0), Cm(0.6),
        size=9, color=GREEN, bold=True)
    txt(sl, "걷기 실천율  ↔  비만율", fx + Cm(0.3), Cm(3.15), Cm(15.0), Cm(0.75),
        size=16, bold=True, color=TEXT)
    txt(sl, "걷기율이 높은 지역일수록 비만율이 뚜렷하게 낮아지는 경향을 확인했습니다.",
        fx + Cm(0.3), Cm(3.95), Cm(15.0), Cm(1.5),
        size=10, color=TEXT, wrap=True)

    # Finding 2
    rect(sl, fx, Cm(8.4), Cm(15.5), Cm(5.5), fill=None,
         line_color=CRIMSON, line_pt=1.5)
    txt(sl, "r = 0.64 / 0.58  (양의 상관)", fx + Cm(0.3), Cm(8.65), Cm(15.0), Cm(0.6),
        size=9, color=CRIMSON, bold=True)
    txt(sl, "비만율  ↔  고혈압/당뇨", fx + Cm(0.3), Cm(9.25), Cm(15.0), Cm(0.75),
        size=16, bold=True, color=TEXT)
    txt(sl, "비만율이 높은 지역에서 고혈압과 당뇨 유병률 또한 동일한 패턴으로 높게 나타납니다.",
        fx + Cm(0.3), Cm(10.05), Cm(15.0), Cm(1.5),
        size=10, color=TEXT, wrap=True)

    page_num(sl, 7)
    return sl


# ── Slide 08 — 다중 회귀 분석: 비만율 영향 요인 ─────────────────────────────────

def s08_regression(prs):
    sl = blank(prs)
    top_bar(sl)
    title_bar(sl, "다중 회귀 분석: 비만율 영향 요인")
    subtitle_line(sl,
        "단순 상관을 넘어, 생활습관 변수들을 동시 통제한 후의 독립적인 영향력을 산출했습니다.")

    # Code block
    rect(sl, Cm(1.3), Cm(2.8), Cm(31.2), Cm(2.0),
         fill=RGBColor(0xF5, 0xF5, 0xFA), line_color=BORDER, line_pt=0.5)
    txt(sl, "R Model Formula", Cm(1.6), Cm(2.95), Cm(30), Cm(0.5),
        size=9, color=TEAL, font=CODE, bold=True)
    txt(sl, "model  <-  lm(비만율  ~  걷기_실천율  +  중강도_활동율  +  고위험_음주율  +  스트레스_인지율  +  흡연율,  data = df)",
        Cm(1.6), Cm(3.5), Cm(30.5), Cm(0.85),
        size=10.5, color=NAVY, font=CODE, wrap=False)

    findings = [
        (BLUE_C,   "FINDING 01", "걷기율의 독립적 효과",
         "음주, 흡연 등 타 변수를 통제한 후에도 걷기 실천율은 비만율 감소에 통계적으로 유의미한 독립적 효과를 미칩니다."),
        (CRIMSON,  "FINDING 02", "고위험 음주율의 위험성",
         "회귀계수 산출 결과, 고위험 음주율은 비만율 증가를 가장 크게 견인하는 치명적인 위험 요인으로 식별되었습니다."),
        (GREEN,    "VALIDATION", "모델 안정성 (VIF 검정)",
         "분산팽창계수(VIF) 검정 결과 모든 변수가 기준치를 통과하여 다중공선성 오류가 없는 안정적 모델임이 확인되었습니다."),
    ]
    fw = Cm(10.2)
    for i, (bc, lbl, ttl, body) in enumerate(findings):
        card(sl, Cm(1.3) + i * (fw + Cm(0.55)), Cm(5.3),
             fw, Cm(7.5), bc, lbl, ttl, body,
             lsz=8, tsz=13.5, bsz=10)

    page_num(sl, 8)
    return sl


# ── Slide 09 — 만성질환 영향 분석 (간접 경로) ────────────────────────────────────

def s09_mechanism(prs):
    sl = blank(prs)
    top_bar(sl)
    title_bar(sl, "만성질환 영향 분석 (간접 경로 규명)")
    subtitle_line(sl,
        "걷기를 많이 하는 상위 25% 지역은 하위 지역에 비해 질환 진단율이 일관되게 낮습니다.")

    # Left text
    rect(sl, Cm(1.3), Cm(3.2), Cm(0.12), Cm(9.5), fill=TEAL)
    txt(sl, "Core Mechanism", Cm(1.7), Cm(3.4), Cm(13.5), Cm(0.55),
        size=9, color=MGRAY)
    txt(sl, "걷기  →  비만 감소  →  만성질환 감소", Cm(1.7), Cm(4.0), Cm(13.5), Cm(0.85),
        size=16, bold=True, color=NAVY)
    txt(sl,
        "고혈압 모델 분석 시 비만율이 1% 증가할 때 고혈압 진단율이 약 0.25%p 증가했습니다.\n\n"
        "걷기의 직접적 효과보다, '비만 감소'를 매개로 한 간접 경로가 만성질환 방어의 핵심 보건 메커니즘으로 작동함을 규명했습니다.",
        Cm(1.7), Cm(5.1), Cm(13.5), Cm(5.5), size=11, color=TEXT, wrap=True)

    # Bar chart (right)
    rect(sl, Cm(16.5), Cm(2.5), Cm(16.0), Cm(11.0),
         fill=None, line_color=BORDER, line_pt=0.5)
    txt(sl, "걷기 수준별 만성질환 진단율 격차", Cm(16.5), Cm(2.7), Cm(16.0), Cm(0.65),
        size=12, bold=True, color=NAVY, align=PP_ALIGN.CENTER)

    chart_data = [
        ("고혈압",  0.68, 0.84),
        ("당뇨",    0.54, 0.71),
    ]
    bar_max_h = Cm(6.0)
    bx_start  = Cm(18.5)
    for ci, (label, hi, lo) in enumerate(chart_data):
        bx = bx_start + ci * Cm(6.0)
        for bar_idx, (val, color, bar_lbl) in enumerate([
            (hi, BLUE_C, "상위 25%"),
            (lo, CRIMSON, "하위 25%"),
        ]):
            bw  = Cm(2.0)
            bh  = bar_max_h * (val / 1.0)
            by  = Cm(12.5) - bh
            bbx = bx + bar_idx * Cm(2.5)
            rect(sl, bbx, by, bw, bh, fill=color)
            txt(sl, bar_lbl, bbx, by - Cm(0.5), bw, Cm(0.45),
                size=8, color=color, align=PP_ALIGN.CENTER, bold=True)
            txt(sl, f"{val:.0%}", bbx, by - Cm(0.95), bw, Cm(0.45),
                size=8.5, color=color, align=PP_ALIGN.CENTER)
        # x-axis label
        txt(sl, label, bx, Cm(12.6), Cm(4.5), Cm(0.6),
            size=11, bold=True, color=TEXT, align=PP_ALIGN.CENTER)

    page_num(sl, 9)
    return sl


# ── Slide 10 — K-Means 군집 분석 ─────────────────────────────────────────────

def s10_kmeans(prs):
    sl = blank(prs)
    top_bar(sl)
    title_bar(sl, "K-Means 군집 분석  ( K = 3 )")
    subtitle_line(sl,
        "변수들을 표준화 후 군집을 적용하여 전국을 3가지 지역 특성 그룹으로 분류했습니다.")

    clusters = [
        (CRIMSON, "⚠", "건강 위험군",
         "비만, 음주, 흡연율 지표가 높게 나타나며 상대적으로 걷기 실천율이 매우 저조한 취약 지역입니다. 강력한 정책 개입이 시급합니다."),
        (AMBER,   "≡", "건강 중간군",
         "전국 평균 수준의 건강 행태와 만성질환 지표를 유지하고 있는 보편적 지역군입니다. 현행 유지 및 점진적 개선이 적합합니다."),
        (GREEN,   "♻", "건강 양호군",
         "걷기 실천율이 활성화되어 있으며, 비만 및 음주 비율이 현저히 낮아 건강 지표가 전반적으로 우수한 모범 지역입니다."),
    ]
    cw = Cm(10.2)
    for i, (bc, icon, title, body) in enumerate(clusters):
        cx = Cm(1.3) + i * (cw + Cm(0.55))
        cy = Cm(2.9)
        ch = Cm(9.0)
        rect(sl, cx, cy, cw, Cm(0.2), fill=bc)
        rect(sl, cx, cy + Cm(0.2), cw, ch - Cm(0.2), fill=None, line_color=BORDER, line_pt=0.5)
        # Icon
        txt(sl, icon, cx, cy + Cm(0.5), cw, Cm(1.0),
            size=22, bold=True, color=bc, align=PP_ALIGN.CENTER)
        txt(sl, title, cx, cy + Cm(1.7), cw, Cm(0.75),
            size=17, bold=True, color=bc, align=PP_ALIGN.CENTER)
        txt(sl, body, cx + Cm(0.3), cy + Cm(2.6), cw - Cm(0.5), ch - Cm(2.8),
            size=10, color=TEXT, wrap=True)

    # Policy implication box
    rect(sl, Cm(1.3), Cm(12.5), Cm(31.2), Cm(1.8),
         fill=RGBColor(0xEE, 0xF5, 0xF5), line_color=TEAL, line_pt=0.5)
    txt(sl, "정책적 시사점",  Cm(1.7), Cm(12.65), Cm(5), Cm(0.55), size=10, bold=True, color=TEAL)
    txt(sl, "공간 지도 시각화 결과, 전국 단위의 획일적 기준이 아닌 지역 유형별 취약 요인을 겨냥한 차별화된 보건 예산 배분이 필수적입니다.",
        Cm(6.5), Cm(12.7), Cm(25.5), Cm(1.4), size=11, color=TEXT, wrap=True)

    page_num(sl, 10)
    return sl


# ── Slide 11 — 머신러닝 비만 예측 변수 중요도 ────────────────────────────────────

def s11_ml(prs):
    sl = blank(prs)
    top_bar(sl)
    title_bar(sl, "머신러닝 비만 예측 변수 중요도")
    subtitle_line(sl,
        "어떤 생활습관이 비만율에 가장 압도적인 영향을 미치는지 Random Forest로 정량화했습니다.")

    # Left description
    txt(sl, "로지스틱 회귀 교차 검증", Cm(1.3), Cm(3.2), Cm(14.5), Cm(0.75),
        size=14, bold=True, color=NAVY)
    txt(sl,
        "가상 시나리오(걷기율 30%, 음주율 20%) 분석 결과, 걷기율이 감소하고 음주율이 증가할수록 비만 고위험군으로 전락할 확률이 급증합니다.\n\n"
        "통계적 방법(다중회귀)과 ML 기법 모두에서 걷기율의 중요성이 일관되게 확인되었습니다.",
        Cm(1.3), Cm(4.15), Cm(14.5), Cm(6.0), size=11, color=TEXT, wrap=True)

    # Feature importance bars (right)
    txt(sl, "변수별 기여도  (ntree=500, IncNodePurity)",
        Cm(17.0), Cm(2.9), Cm(15.5), Cm(0.65), size=11, bold=True, color=NAVY)

    features = [
        ("흡연율",       0.25, MGRAY),
        ("고위험 음주율", 0.22, CRIMSON),
        ("스트레스 인지율", 0.19, MGRAY),
        ("중강도 신체활동",0.18, MGRAY),
        ("걷기 실천율",  0.16, TEAL),
    ]
    max_bar_w = Cm(10.0)
    for i, (label, val, color) in enumerate(features):
        fy = Cm(3.9) + i * Cm(1.8)
        txt(sl, label, Cm(17.0), fy, Cm(4.8), Cm(0.65), size=11, color=TEXT)
        bh = Cm(0.65)
        bw = max_bar_w * val
        rect(sl, Cm(22.0), fy + Cm(0.02), bw, bh, fill=color)
        rect(sl, Cm(22.0) + bw, fy + Cm(0.02), max_bar_w - bw, bh,
             fill=LGRAY)
        txt(sl, f"~{val:.0%}", Cm(22.0) + max_bar_w + Cm(0.2), fy, Cm(1.8), Cm(0.65),
            size=11, bold=True, color=color)

    page_num(sl, 11)
    return sl


# ── Slide 12 — 패널 회귀 분석 (Fixed Effects) ─────────────────────────────────

def s12_panel(prs):
    sl = blank(prs)
    top_bar(sl)
    title_bar(sl, "패널 회귀 분석 (Fixed Effects)")
    subtitle_line(sl,
        "지역 고유의 인프라와 지형 특성을 완전 통제하여 순수한 걷기율의 인과 효과를 입증합니다.")

    # Left
    rect(sl, Cm(1.3), Cm(3.2), Cm(0.12), Cm(8.5), fill=TEAL)
    txt(sl, "왜 패널 모델인가?", Cm(1.7), Cm(3.2), Cm(13.5), Cm(0.75),
        size=14, bold=True, color=NAVY)
    txt(sl,
        "plm 패키지의 Fixed Effects 모델을 사용하여 특정 지역이 가진 고유의 인프라 차이를 통제했습니다.\n\n"
        "분석 결과 걷기율의 Estimate가 음수(−),  p < 0.05 로 유의하게 산출되어, 걷기 정책이 전국 단위로 확고한 효과가 있음을 증명합니다.",
        Cm(1.7), Cm(4.2), Cm(13.5), Cm(7.0), size=11, color=TEXT, wrap=True)

    # Right: trend chart
    rect(sl, Cm(16.5), Cm(2.5), Cm(16.0), Cm(11.0),
         fill=None, line_color=BORDER, line_pt=0.5)
    txt(sl, "흡연율 vs 체중조절 시도율 교차 트렌드 (2009-2030)",
        Cm(16.5), Cm(2.7), Cm(16.0), Cm(0.65),
        size=11, bold=True, color=NAVY, align=PP_ALIGN.CENTER)

    # X axis labels
    txt(sl, "2009", Cm(17.0), Cm(13.0), Cm(2.5), Cm(0.5), size=9, color=MGRAY)
    txt(sl, "2030(E)", Cm(29.5), Cm(13.0), Cm(2.5), Cm(0.5), size=9, color=MGRAY)

    # Declining line (흡연율) — using rect as line segment approximation
    # Line 1: 체중조절 의지 (rising, blue) y: 11→4 over x: 17.5→31
    from pptx.util import Emu as _Emu
    # Draw crossing lines with thin rects at angle (approximate with two rects)
    # Rising: body weight control
    rect(sl, Cm(17.5), Cm(9.5), Cm(13.0), Cm(0.12), fill=BLUE_C)   # approx rising
    rect(sl, Cm(17.5), Cm(9.5), Cm(0.12), Cm(3.0), fill=BLUE_C)    # left start
    rect(sl, Cm(30.5), Cm(6.5), Cm(0.12), Cm(3.0), fill=BLUE_C)    # right end
    txt(sl, "체중조절 의지 ↑", Cm(28.5), Cm(6.0), Cm(4.0), Cm(0.55),
        size=9, bold=True, color=BLUE_C)

    # Falling: smoking rate (crimson)
    rect(sl, Cm(17.5), Cm(6.5), Cm(13.0), Cm(0.12), fill=CRIMSON)   # approx falling
    rect(sl, Cm(17.5), Cm(6.5), Cm(0.12), Cm(3.0), fill=CRIMSON)
    rect(sl, Cm(30.5), Cm(9.5), Cm(0.12), Cm(3.0), fill=CRIMSON)
    txt(sl, "흡연율 ↓", Cm(29.0), Cm(10.0), Cm(3.0), Cm(0.55),
        size=9, bold=True, color=CRIMSON)

    page_num(sl, 12)
    return sl


# ── Slide 13 — 정책 시뮬레이션 (2026-2030 예측) ──────────────────────────────────

def s13_simulation(prs):
    sl = blank(prs)
    top_bar(sl)
    title_bar(sl, "정책 시뮬레이션 (2026-2030 예측)")
    subtitle_line(sl,
        "걷기율을 인위적으로 높이면 향후 5년 내 비만율의 가파른 상승세를 방어할 수 있는가?")

    scenarios = [
        (CRIMSON, "현상 유지 시나리오",
         "과거 추세가 그대로 연장되어 비만율이 2030년까지 브레이크 없이 지속적으로 가파르게 상승합니다."),
        (GREEN,   "정책 개입 시나리오",
         "매년 걷기율을 +1~5%p 상향 시킬 경우, 2030년 기준 현상 유지 대비 약 3.5%p의 비만율 상승 억제 효과를 가져옵니다."),
    ]
    for i, (bc, ttl, body) in enumerate(scenarios):
        sy = Cm(2.8) + i * Cm(5.5)
        rect(sl, Cm(1.3), sy, Cm(14.5), Cm(0.18), fill=bc)
        rect(sl, Cm(1.3), sy + Cm(0.18), Cm(14.5), Cm(4.8),
             fill=None, line_color=BORDER, line_pt=0.5)
        txt(sl, ttl,  Cm(1.6), sy + Cm(0.45), Cm(14.0), Cm(0.75), size=14, bold=True, color=bc)
        txt(sl, body, Cm(1.6), sy + Cm(1.3),  Cm(14.0), Cm(3.5), size=11, color=TEXT, wrap=True)

    # Forecast chart (right)
    rect(sl, Cm(16.5), Cm(2.5), Cm(16.0), Cm(11.0),
         fill=None, line_color=BORDER, line_pt=0.5)
    txt(sl, "2030 비만율 예측 다이어그램",
        Cm(16.5), Cm(2.7), Cm(16.0), Cm(0.65),
        size=12, bold=True, color=NAVY, align=PP_ALIGN.CENTER)

    txt(sl, "2025 기준점", Cm(18.0), Cm(9.2), Cm(4.0), Cm(0.55), size=9, color=MGRAY)
    # Base point
    rect(sl, Cm(20.5), Cm(9.0), Cm(0.4), Cm(0.4), fill=NAVY)
    # Rising line (현상 유지)
    rect(sl, Cm(20.5), Cm(6.5), Cm(9.5), Cm(0.12), fill=CRIMSON)
    rect(sl, Cm(20.5), Cm(6.5), Cm(0.12), Cm(2.5), fill=CRIMSON)
    rect(sl, Cm(30.0), Cm(6.5), Cm(0.12), Cm(0.8), fill=CRIMSON)
    txt(sl, "현상 유지", Cm(28.5), Cm(5.9), Cm(3.5), Cm(0.55), size=9, bold=True, color=CRIMSON)
    # Slower line (정책 개입)
    rect(sl, Cm(20.5), Cm(8.0), Cm(9.5), Cm(0.12), fill=GREEN)
    rect(sl, Cm(20.5), Cm(8.0), Cm(0.12), Cm(1.0), fill=GREEN)
    rect(sl, Cm(30.0), Cm(8.0), Cm(0.12), Cm(0.8), fill=GREEN)
    txt(sl, "정책 개입", Cm(28.5), Cm(7.8), Cm(3.5), Cm(0.55), size=9, bold=True, color=GREEN)
    # Gap label
    txt(sl, "~3.5%p 방어", Cm(25.5), Cm(7.0), Cm(4.0), Cm(0.55),
        size=10, bold=True, color=NAVY)

    txt(sl, "2025", Cm(17.5), Cm(12.8), Cm(2.5), Cm(0.5), size=9, color=MGRAY)
    txt(sl, "2030(E)", Cm(29.0), Cm(12.8), Cm(2.8), Cm(0.5), size=9, color=MGRAY)

    page_num(sl, 13)
    return sl


# ── Slide 14 — 4대 건강지표 통합 미래 트렌드 ──────────────────────────────────────

def s14_trends(prs):
    sl = blank(prs)
    top_bar(sl)
    title_bar(sl, "4대 건강지표 통합 미래 트렌드")
    subtitle_line(sl,
        "과거의 궤적을 통해 도출한 2030년 질환 및 건강행태의 거시적 변화 4대 추이입니다.")

    trends = [
        (CRIMSON,  "TREND 01", "비만율 & 당뇨 동반 상승",
         "비만율 상승 궤적에 따라 당뇨 진단율이 시차를 두고 비례하여 동반 상승하는 패턴이 명확히 예측됩니다."),
        (BLUE_C,   "TREND 02", "고혈압 & 당뇨 상호 강화",
         "강력한 양의 상관관계를 가진 고혈압과 당뇨 지표는 2030년까지 상호 위험을 강화하며 동반 확산될 전망입니다."),
        (AMBER,    "TREND 03", "체중조절 의지 vs 비만 역설",
         "시민들의 체중조절 시도율(관심)은 지속 상승하나 실제 비만율도 상승하는 역설적 현상입니다. 실천 환경 조성이 급선무입니다."),
        (TEAL,     "TREND 04", "비만 확산과 고혈압 경고",
         "비만율 누적 상승치가 임계점을 넘으며, 향후 5년간 고혈압 진단율의 상승폭 또한 매우 가팔라질 것으로 경고됩니다."),
    ]

    tw = Cm(15.3)
    th = Cm(6.2)
    positions = [(Cm(1.3), Cm(2.9)), (Cm(17.2), Cm(2.9)),
                 (Cm(1.3), Cm(9.7)), (Cm(17.2), Cm(9.7))]
    for (tx, ty), (bc, lbl, ttl, body) in zip(positions, trends):
        card(sl, tx, ty, tw, th, bc, lbl, ttl, body,
             lsz=8, tsz=14, bsz=10.5)

    page_num(sl, 14)
    return sl


# ── Slide 15 — 최종 분석 결론 및 정책 제언 ─────────────────────────────────────

def s15_conclusion(prs):
    sl = blank(prs)
    # Dark background (top portion)
    rect(sl, Cm(0), Cm(0), W, Cm(2.2), fill=DKBG)
    txt(sl, "15.  Conclusion", Cm(0), Cm(0.4), W, Cm(0.65),
        size=11, color=MGRAY, align=PP_ALIGN.CENTER)
    txt(sl, "최종 분석 결론 및 정책 제언",
        Cm(0), Cm(1.0), W, Cm(0.95),
        size=22, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    # Findings box
    rect(sl, Cm(1.5), Cm(2.7), Cm(30.8), Cm(8.2),
         fill=LGRAY, line_color=BORDER, line_pt=0.5)

    conclusions = [
        ("1. 유의성 검증:",
         "걷기 실천율 향상은 비만율 감소에 기여함이 회귀분석 및 패널(FE) 모델 모두에서 유의하게 확인되었습니다."),
        ("2. 경로 규명:",
         "단순 걷기 효과를 넘어 '걷기 → 비만 감소 → 만성질환 감소'로 이어지는 간접 경로 메커니즘을 밝혔습니다."),
        ("3. 지역 격차 심각:",
         "서울(68%) vs 강원(38%) 등 인프라 불평등이 심각하여 차별화된 타겟팅 예산 배분이 요구됩니다."),
        ("4. 복합 정책 필수:",
         "걷기 인프라 확충 단독 개입의 한계를 넘어, 금연과 절주를 결합한 시너지 보건 정책 설계가 필수적입니다."),
    ]

    for i, (bold_part, rest) in enumerate(conclusions):
        cy = Cm(3.1) + i * Cm(1.9)
        txt(sl, bold_part, Cm(2.0), cy, Cm(7.0), Cm(0.75),
            size=12, bold=True, color=NAVY)
        txt(sl, rest, Cm(9.0), cy, Cm(22.5), Cm(1.6),
            size=11, color=TEXT, wrap=True)

    # Closing quote
    rect(sl, Cm(1.5), Cm(11.5), Cm(30.8), Cm(4.5), fill=DKBG)
    txt(sl,
        '"229개 시군구 단위 걷기 실천율과 만성질환의 정량적 관계 연구는\n'
        '다층적 통계 방법론을 통해 실효성 있는 보건 정책의 핵심 근거를 구축했음을 증명합니다."',
        Cm(2.5), Cm(12.3), Cm(28.8), Cm(2.8),
        size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER, wrap=True)

    page_num(sl, 15)
    return sl


# ── Main ───────────────────────────────────────────────────────────────────────

def main():
    prs = new_prs()

    print("슬라이드 생성 중...")
    funcs = [
        s01_title,
        s02_background,
        s03_toc,
        s04_data,
        s05_architecture,
        s06_preprocessing,
        s07_heatmap,
        s08_regression,
        s09_mechanism,
        s10_kmeans,
        s11_ml,
        s12_panel,
        s13_simulation,
        s14_trends,
        s15_conclusion,
    ]
    for i, fn in enumerate(funcs, 1):
        fn(prs)
        print(f"  [{i:02d}/15] {fn.__name__} 완료")

    prs.save(str(OUTPUT))
    size_mb = OUTPUT.stat().st_size / 1024 / 1024
    print(f"\n저장 완료: {OUTPUT.name}  ({size_mb:.1f} MB, 15 슬라이드)")


if __name__ == "__main__":
    main()
